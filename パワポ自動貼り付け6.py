import os
from pptx import Presentation
from pptx.util import Inches

# 画像が保存されているフォルダを指定
image_dir = r"C:\Users\nmukaida\Desktop\0627ボックスカルバート"

# 新しいPowerPointプレゼンテーションを作成
prs = Presentation()

# フォルダ内の全ての.pngファイルに対してループを行う
for filename in os.listdir(image_dir):
    if filename.endswith(".png"):
        # 新しいスライドを作成（レイアウトは「タイトルとコンテンツ」）
        slide_layout = prs.slide_layouts[1] 
        slide = prs.slides.add_slide(slide_layout)

        # タイトルにファイル名を設定
        title = slide.shapes.title
        title.text = os.path.splitext(filename)[0]
        
        # コンテンツ部分に画像を設定
        left = Inches(0.79)
        top = Inches(1.8)
        width = Inches(8.27)
        height = Inches(4.92)
        slide.shapes.add_picture(os.path.join(image_dir, filename), left, top, width, height)


# プレゼンテーションを保存
prs.save("presentation.pptx")
print("プレゼンテーションが保存されました")